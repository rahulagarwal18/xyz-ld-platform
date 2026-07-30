import collections
import collections.abc
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_eyecatchy_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Color Palette Definitions ──
    COLOR_BG = RGBColor(248, 250, 252)       # Slate 50
    COLOR_CARD = RGBColor(255, 255, 255)     # White card
    COLOR_BORDER = RGBColor(226, 232, 240)   # Slate 200
    COLOR_NAVY = RGBColor(0, 32, 91)         # xyz Navy
    COLOR_ACCENT = RGBColor(0, 156, 222)     # xyz Blue
    COLOR_RED = RGBColor(204, 0, 0)          # Alert Red
    COLOR_TEXT = RGBColor(30, 41, 59)        # Slate 800
    COLOR_MUTED = RGBColor(100, 116, 139)    # Slate 500
    COLOR_WHITE = RGBColor(255, 255, 255)

    FONT_FAMILY = "Segoe UI"

    # Image Paths
    IMG_DIR = "B:\\AI_Projects\\xyz-ld-platform\\public\\images"
    img_ai = os.path.join(IMG_DIR, "ai_mindset.jpg")
    img_career = os.path.join(IMG_DIR, "career_xyz.jpg")
    img_leading = os.path.join(IMG_DIR, "leading_xyz.jpg")

    # Helper: Adds base slide with clean typography and bottom accent footer
    def add_base_slide(title_text):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Full page background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

        # Modern Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY

        # Accent Line under Title
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.fill.background()

        # Footer branding
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3))
        ftf = footer_box.text_frame
        ftf.word_wrap = True
        ftf.margin_left = ftf.margin_top = ftf.margin_bottom = ftf.margin_right = 0
        fp = ftf.paragraphs[0]
        fp.text = "xyz Learning & Development Department  •  TLCE LMS Portfolio"
        fp.font.name = FONT_FAMILY
        fp.font.size = Pt(9)
        fp.font.bold = True
        fp.font.color.rgb = COLOR_MUTED

        return slide

    # Helper: Adds a card
    def add_card(slide, left, top, width, height, border_color=COLOR_BORDER, fill_color=COLOR_CARD):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # Helper: Writes text inside a card
    def write_card_content(slide, card_left, card_top, card_width, card_height, heading, items, is_alert=False, text_color=COLOR_TEXT):
        pad = Inches(0.25)
        tx_box = slide.shapes.add_textbox(card_left + pad, card_top + pad, card_width - (pad * 2), card_height - (pad * 2))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0

        p = tf.paragraphs[0]
        p.text = heading
        p.font.name = FONT_FAMILY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_RED if is_alert else COLOR_NAVY
        p.space_after = Pt(12)

        for title, desc in items:
            p_title = tf.add_paragraph()
            p_title.text = "•  " + title
            p_title.font.name = FONT_FAMILY
            p_title.font.size = Pt(11)
            p_title.font.bold = True
            p_title.font.color.rgb = text_color
            p_title.space_before = Pt(6)

            p_desc = tf.add_paragraph()
            p_desc.text = "    " + desc
            p_desc.font.name = FONT_FAMILY
            p_desc.font.size = Pt(10)
            p_desc.font.color.rgb = COLOR_MUTED
            p_desc.space_after = Pt(2)

    # ──────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide (Sleek Dark Accent Block)
    # ──────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_WHITE
    bg.line.fill.background()

    # Dark blue block on the left
    left_block = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = COLOR_NAVY
    left_block.line.fill.background()

    # Image overlay in the left block with transparency-like feel
    if os.path.exists(img_ai):
        pic = s1.shapes.add_picture(img_ai, Inches(0.2), Inches(1.8), Inches(4.1), Inches(3.2))

    brand_box = s1.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(3), Inches(1.5))
    btf = brand_box.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.text = "xyz"
    bp.font.name = FONT_FAMILY
    bp.font.size = Pt(54)
    bp.font.bold = True
    bp.font.color.rgb = COLOR_WHITE
    
    bp_sub = btf.add_paragraph()
    bp_sub.text = "LEARNING & DEVELOPMENT"
    bp_sub.font.name = FONT_FAMILY
    bp_sub.font.size = Pt(11)
    bp_sub.font.bold = True
    bp_sub.font.color.rgb = COLOR_ACCENT
    bp_sub.space_before = Pt(4)

    title_box = s1.shapes.add_textbox(Inches(5.3), Inches(2.2), Inches(7.2), Inches(4.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "TLCE LMS Portal"
    p1.font.name = FONT_FAMILY
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_NAVY
    p1.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "Modernized Training, Learning, Conference & Engagement Platform"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ACCENT
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "A comprehensive suite built with strict corporate color alignment, anti-cheat assessment tracking, automated notifications, and real-time manager notification routing."
    p3.font.name = FONT_FAMILY
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_after = Pt(24)

    p4 = tf.add_paragraph()
    p4.text = "Project Release  •  Version 2.0"
    p4.font.name = FONT_FAMILY
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_TEXT

    # ──────────────────────────────────────────────────────────
    # SLIDE 2: Brand Identity (Hero Image Split Layout)
    # ──────────────────────────────────────────────────────────
    s2 = add_base_slide("Executive Summary & Brand Identity")
    
    # Left Side: Hero Image Showcase
    if os.path.exists(img_career):
        img_border = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
        img_border.fill.solid()
        img_border.fill.fore_color.rgb = COLOR_BORDER
        img_border.line.fill.background()
        # Actual Image overlay
        s2.shapes.add_picture(img_career, Inches(0.9), Inches(1.7), Inches(5.4), Inches(4.8))

    # Right Side: Card details
    add_card(s2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0))
    write_card_content(s2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), 
        "Corporate Color & Identity", [
            ("Platform Mission", "Designed to deliver and coordinate training modules for xyz engineering, design, strategy and product teams."),
            ("Corporate Navy (#003087)", "Central theme accent used for headers, primary buttons, and tabs."),
            ("Accent Blue (#009CDE)", "Interactive highlights applied to checked-in lists, calendar badging, and progress sliders."),
            ("White Dashboard Layout", "Minimalist white pages replace standard styling configurations, providing readability and corporate formatting.")
        ]
    )

    # ──────────────────────────────────────────────────────────
    # SLIDE 3: Yearly Calendar Grid & Catalog (3-Card Grid)
    # ──────────────────────────────────────────────────────────
    s3 = add_base_slide("Program Catalog & Yearly Calendar")
    
    card_width = Inches(3.64)
    # Card 1: Banner Carousel
    add_card(s3, Inches(0.8), Inches(1.6), card_width, Inches(5.0))
    write_card_content(s3, Inches(0.8), Inches(1.6), card_width, Inches(5.0), 
        "1. Immersive Carousel", [
            ("Auto-Play Slider", "Featured banner program showing key upcoming courses."),
            ("Details Overlay", "Integrates dates, times, key objectives, and registration triggers on the slider."),
            ("Enrolled Tags", "Displays enrolled status directly on slides.")
        ]
    )

    # Card 2: Program Roster
    add_card(s3, Inches(4.84), Inches(1.6), card_width, Inches(5.0))
    write_card_content(s3, Inches(4.84), Inches(1.6), card_width, Inches(5.0), 
        "2. Roster Grid Cards", [
            ("Metadata Badging", "Displays duration hours, target audiences, and meal tags."),
            ("Dynamic Capacity Status", "Color progress bars transition from green (available) to red (fully booked)."),
            ("Direct CTA Triggers", "Launches registration modals immediately.")
        ]
    )

    # Card 3: 12-Month Calendar
    add_card(s3, Inches(8.88), Inches(1.6), card_width, Inches(5.0))
    write_card_content(s3, Inches(8.88), Inches(1.6), card_width, Inches(5.0), 
        "3. Yearly Calendar Grid", [
            ("12-Month Flow", "Includes one highlighted program scheduled for each month of the year."),
            ("Key Focus Highlights", "Displays summaries of monthly objectives (e.g. Design Systems, Strategic Business Acumen)."),
            ("Status Indicators", "Updates calendar items dynamically on registration.")
        ]
    )

    # ──────────────────────────────────────────────────────────
    # SLIDE 4: Registration, Meals & Cabs (Visual Split)
    # ──────────────────────────────────────────────────────────
    s4 = add_base_slide("Registration, Meals & Logistics")
    
    # Left Card
    add_card(s4, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    write_card_content(s4, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), 
        "Registration & Seat Allocation", [
            ("Over-Capacity Waitlisting", "Automatically disables registration button and displays 'Join Waitlist' if seats hit 100% capacity."),
            ("Veg/Non-Veg Meal Options", "Complimentary options restricted strictly to 'Veg' and 'Non-Veg'. Only displays on events providing lunch/dinner."),
            ("User Dashboard Badging", "Highlights confirmations and waitlist order indices.")
        ]
    )

    # Right Card: Cab Shuttle
    add_card(s4, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0))
    write_card_content(s4, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), 
        "Corporate Shuttle booking System", [
            ("Standalone Toggle Switch", "Replaced standard label+input toggle with a custom div component to prevent double-fire toggle errors."),
            ("Drop Location Stations", "Selectable pickup zones (e.g. City Tech Park, Tower A Plaza, East Business Bay)."),
            ("Timetable Batches", "Accommodates specific shift pick-ups (08:30 AM, 09:15 AM, 10:00 AM, 01:15 PM)."),
            ("Visual Confirmation Tag", "Generates summary details dynamically on the modal.")
        ]
    )

    # ──────────────────────────────────────────────────────────
    # SLIDE 5: Scheduled Timelines (Flow Template)
    # ──────────────────────────────────────────────────────────
    s5 = add_base_slide("Automated Email Timeline triggers")

    # Draw horizontal timeline flow card
    add_card(s5, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.1))

    # Intro text
    introBox = s5.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(0.8))
    itf = introBox.text_frame
    itf.word_wrap = True
    p = itf.paragraphs[0]
    p.text = "Timeline-Triggered Reminders & Logs"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(4)
    
    p_sub = itf.add_paragraph()
    p_sub.text = "Automated email reminders are dispatched to the simulated Inbox Logs view based on program date parameters:"
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(11)
    p_sub.font.color.rgb = COLOR_MUTED

    # Timeline cards
    step_width = Inches(3.4)
    step_height = Inches(3.2)
    step_top = Inches(2.8)
    
    steps = [
        ("Step 1: 2 Days Prior", "Dispatched to confirmation registries", "Contains logistical coordinates, map pins, venue location, and shuttle cab options."),
        ("Step 2: 1 Day Prior", "High priority cert reminder", "Urges participants to complete mandatory pre-program assessments in fullscreen mode."),
        ("Step 3: Day-Of (6 AM)", "Immediate session trigger", "Delivers final stream login credentials and pick-up details to the attendee dashboard.")
    ]

    for idx, (step_title, step_sub, step_desc) in enumerate(steps):
        s_left = Inches(1.2) + (idx * Inches(3.75))
        # Draw step card
        scard = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s_left, step_top, step_width, step_height)
        scard.fill.solid()
        scard.fill.fore_color.rgb = COLOR_BG
        scard.line.color.rgb = COLOR_BORDER
        scard.line.width = Pt(1.5)

        # Draw step text
        s_box = s5.shapes.add_textbox(s_left + Inches(0.2), step_top + Inches(0.2), step_width - Inches(0.4), step_height - Inches(0.4))
        stf = s_box.text_frame
        stf.word_wrap = True
        
        p = stf.paragraphs[0]
        p.text = step_title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        p.space_after = Pt(4)

        p_sub = stf.add_paragraph()
        p_sub.text = step_sub
        p_sub.font.name = FONT_FAMILY
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_ACCENT
        p_sub.space_after = Pt(12)

        p_desc = stf.add_paragraph()
        p_desc.text = step_desc
        p_desc.font.name = FONT_FAMILY
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = COLOR_TEXT

    # ──────────────────────────────────────────────────────────
    # SLIDE 6: Assessment Suite & Anti-Cheat (Hero Alert Split Layout)
    # ──────────────────────────────────────────────────────────
    s6 = add_base_slide("Assessment Suite & Exam Integrity")
    
    # Left Card: Assessments
    add_card(s6, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    write_card_content(s6, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), 
        "Assessment Delivery Engine", [
            ("Integrated Pre/Post Tests", "Supports Multiple Choice Questions (MCQ) and One-Word short answer text fields."),
            ("Dynamic Question Injection", "Loads category-specific question sets (Engineering, Strategy, UX) based on the active event category."),
            ("Copy-Paste Blockers", "Blocks keyboard copy/paste commands (Ctrl+C, Ctrl+V), right-click context menu, and text drag-selections."),
            ("Immediate Roster Recording", "Scores >= 60% are marked 'Passed' immediately on the admin dashboard console.")
        ]
    )

    # Right Card: Fullscreen Cheat Warning (Red Alert Border Card)
    add_card(s6, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), border_color=COLOR_RED)
    write_card_content(s6, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), 
        "Fullscreen Anti-Cheat Mechanism", [
            ("Mandatory Fullscreen Mode", "Tests launch in mandatory fullscreen. The browser event listener actively monitors window focus parameters."),
            ("3-Strike Exit Rule", "Exiting fullscreen prompts a strike warning. On the 3rd strike, the test auto-terminates."),
            ("Auto-Failed Status", " disallows further attempts. The exam is marked 'Failed (Cheating)'.")
        ], is_alert=True
    )

    # ──────────────────────────────────────────────────────────
    # SLIDE 7: Feedback Engine & L&D AI Assistant
    # ──────────────────────────────────────────────────────────
    s7 = add_base_slide("Feedback Engine & L&D AI Assistant")
    
    add_card(s7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    write_card_content(s7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), 
        "Feedback & Evaluation System", [
            ("4-Category Rating Metrics", "Enables 5-star evaluation of Overall Experience, Speaker Quality, Content Relevance, and Logistics."),
            ("Detailed Comment Feedbacks", "Includes text boxes for participants to explain practical application and improvement feedback."),
            ("Admin Log Dispatches", "Submitting feedback broadcasts full response logs directly to the Admin Inbox Drawer.")
        ]
    )

    add_card(s7, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0))
    write_card_content(s7, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), 
        "Compact AI Chat Assistant", [
            ("Floating Picture-in-Picture", "Launches a compact 360x500px overlay in the bottom-right corner rather than a full page block."),
            ("Quick FAQ Selection Chips", "Includes 6 preset question chips (e.g. 'Cab pickup options?') for instant queries without typing."),
            ("Groq Llama 3.3 Integration", "Provides prompt responses. System instructions block internal model name leakage.")
        ]
    )

    # ──────────────────────────────────────────────────────────
    # SLIDE 8: Analytics & Learning Hours (Dashboard Table Card)
    # ──────────────────────────────────────────────────────────
    s8 = add_base_slide("Total Learning Hours & Analytics Dashboard")
    
    # Text intro on a clean single wide card
    add_card(s8, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.1))
    
    # Description text inside card
    tx_box = s8.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.133), Inches(1.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = "Analytics Calculations & Formulas"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(6)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Total Learning Hours = Σ (Program Duration in Hours × Registered Participants) across all 12 annual courses. The table below represents live analytics metrics generated by the platform:"
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(11)
    p_sub.font.color.rgb = COLOR_MUTED

    # Add a styled table
    rows = 5
    cols = 4
    left = Inches(1.1)
    top = Inches(2.8)
    width = Inches(11.133)
    height = Inches(3.2)
    
    table_shape = s8.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    table.columns[0].width = Inches(4.533)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(2.2)
    
    headers = ["Program Name (Month)", "Duration (Hours)", "Registered Attendees", "Learning Hours Produced"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_FAMILY
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
    data = [
        ["Strategic Thinking & Business Acumen (Jan)", "4 hours", "20 participants", "80 hours"],
        ["Generative AI for Enterprise Workflows (Apr)", "4 hours", "32 participants", "128 hours"],
        ["Leading at xyz (May)", "3 hours", "60 participants (FULL)", "180 hours"],
        ["Design Systems & Enterprise UX (Nov)", "6 hours", "22 participants", "132 hours"]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = COLOR_BG
            else:
                cell.fill.fore_color.rgb = COLOR_WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_FAMILY
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT
            p.alignment = PP_ALIGN.CENTER

    # ──────────────────────────────────────────────────────────
    # SLIDE 9: Technology Stack & Architecture (3-Column Layout)
    # ──────────────────────────────────────────────────────────
    s9 = add_base_slide("Technology Stack & Architecture")
    
    col_width = Inches(3.64)
    # React column
    add_card(s9, Inches(0.8), Inches(1.6), col_width, Inches(5.0))
    write_card_content(s9, Inches(0.8), Inches(1.6), col_width, Inches(5.0), 
        "Frontend Technology Layer", [
            ("React 18 & Vite", "Ultra-fast hot module reloading and UI response times."),
            ("Lucide Icons Library", "Consistent vector icons mapped to action tags."),
            ("Inline Variable CSS", "Strict colors mapped to pre-defined corporate themes.")
        ]
    )

    # Context column
    add_card(s9, Inches(4.84), Inches(1.6), col_width, Inches(5.0))
    write_card_content(s9, Inches(4.84), Inches(1.6), col_width, Inches(5.0), 
        "State Engine Layer", [
            ("React Context API", "Tracks registrations, waitlists, feedbacks, and test metrics."),
            ("Dynamic State Hooks", "Provides centralized data access across tabs.")
        ]
    )

    # LocalStorage column
    add_card(s9, Inches(8.88), Inches(1.6), col_width, Inches(5.0))
    write_card_content(s9, Inches(8.88), Inches(1.6), col_width, Inches(5.0), 
        "Persistence & Emails", [
            ("LocalStorage Hooks", "Ensures check-ins, notifications, and results persist across user browser reboots."),
            ("Simulated Email Engine", "Logs email dispatches directly to the Admin Inbox Drawer.")
        ]
    )

    # ──────────────────────────────────────────────────────────
    # SLIDE 10: Conclusion & Git Deployment Success (Image Split)
    # ──────────────────────────────────────────────────────────
    s10 = add_base_slide("Conclusion & Deployment Summary")
    
    # Left Card
    add_card(s10, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    write_card_content(s10, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), 
        "Deployment Release Success", [
            ("Git Repository Synced", "All modifications committed and pushed successfully to the main branch on the origin remote repository."),
            ("Vercel Live Build", "Auto-compiled and published to Vercel production edge servers, live and fully active."),
            ("Roster Workflows Verified", "Full path testing (login -> check in -> assessment validation -> feedback dispatch) verified successfully.")
        ]
    )

    # Right Side: Image
    if os.path.exists(img_leading):
        img_border = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0))
        img_border.fill.solid()
        img_border.fill.fore_color.rgb = COLOR_BORDER
        img_border.line.fill.background()
        s10.shapes.add_picture(img_leading, Inches(7.0), Inches(1.7), Inches(5.4), Inches(4.8))

    # Save v3
    prs.save('B:\\AI_Projects\\xyz-ld-platform\\xyz_TLCE_LMS_Presentation_v3.pptx')
    print("Premium presentation created successfully as B:\\AI_Projects\\xyz-ld-platform\\xyz_TLCE_LMS_Presentation_v3.pptx")

if __name__ == '__main__':
    create_eyecatchy_deck()
